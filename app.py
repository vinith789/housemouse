from flask import Flask, request, render_template, redirect, url_for, send_from_directory
import os
from werkzeug.utils import secure_filename
import fitz  # PyMuPDF for PDF processing
from pptx import Presentation
from PIL import Image, ImageDraw

import subprocess

UPLOAD_FOLDER = 'uploads'
IMG_FOLDER = 'pre_img'
ALLOWED_EXTENSIONS = {'pdf', 'pptx'}

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(IMG_FOLDER, exist_ok=True)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/')
def index():
    return render_template('upload.html')

@app.route('/upload', methods=['POST'])
def upload():
    file = request.files.get('file')
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)

        # Clear previous images
        for f in os.listdir(IMG_FOLDER):
            os.remove(os.path.join(IMG_FOLDER, f))

        try:
            if filename.endswith('.pdf'):
                convert_pdf_to_images(filepath)
            elif filename.endswith('.pptx'):
                convert_ppt_to_images(filepath)
            return redirect(url_for('viewer'))
        except Exception as e:
            return f"File processing failed: {e}", 500
    return "Invalid file type or upload failed!", 400

def convert_pdf_to_images(pdf_path):
    """Convert PDF to images using PyMuPDF (fitz)."""
    pdf_document = fitz.open(pdf_path)
    for page_number in range(len(pdf_document)):
        page = pdf_document.load_page(page_number)
        pix = page.get_pixmap()
        output_path = os.path.join(IMG_FOLDER, f'slide_{page_number}.jpg')
        pix.save(output_path)
    pdf_document.close()

def convert_ppt_to_images(ppt_path):
    """Convert PowerPoint slides to images using python-pptx and Pillow."""
    prs = Presentation(ppt_path)
    for i, slide in enumerate(prs.slides):
        # Create a blank image with white background
        width = int(prs.slide_width * 0.75)  # Convert EMU to pixels
        height = int(prs.slide_height * 0.75)
        img = Image.new('RGB', (width, height), color='white')
        draw = ImageDraw.Draw(img)

        # Render slide content (basic text rendering as an example)
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                draw.text((10, 10 + i * 20), text, fill='black')

        img.save(os.path.join(IMG_FOLDER, f'slide_{i}.jpg'))

@app.route('/viewer')
def viewer():
    images = [f for f in os.listdir(IMG_FOLDER) if f.endswith('.jpg')]
    images.sort()  # Ensure images are in order
    return render_template('start_viewer.html', images=images)

@app.route('/images/<filename>')
def serve_image(filename):
    return send_from_directory(IMG_FOLDER, filename)

@app.route('/run_viewer')
def run_viewer_route():
    try:
        # Run viewer.py as a subprocess
        result = subprocess.run(['python', 'viewer.py'], capture_output=True, text=True, cwd=os.path.dirname(__file__))
        return f"Viewer executed successfully! Output: <pre>{result.stdout}</pre>"
    except Exception as e:
        return f"Error executing viewer.py: {e}"

if __name__ == '__main__':
    app.run(debug=True)